from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "tools" / "pptx_deps"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


OUT = ROOT / "docs" / "fitness_gym_defense_presentation_pretty.pptx"

COLORS = {
    "ink": RGBColor(15, 23, 42),
    "muted": RGBColor(71, 85, 105),
    "blue": RGBColor(37, 99, 235),
    "green": RGBColor(22, 163, 74),
    "amber": RGBColor(245, 158, 11),
    "red": RGBColor(220, 38, 38),
    "cyan": RGBColor(8, 145, 178),
    "bg": RGBColor(248, 250, 252),
    "card": RGBColor(255, 255, 255),
    "line": RGBColor(226, 232, 240),
    "dark": RGBColor(16, 42, 67),
}


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_fill(bg, COLORS["bg"])
    bg.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.62))
    set_fill(bar, COLORS["dark"])
    bar.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.62), Inches(0.12), Inches(6.88))
    set_fill(accent, COLORS["green"])
    accent.line.fill.background()


def add_title(slide, title, number):
    box = slide.shapes.add_textbox(Inches(0.42), Inches(0.12), Inches(11.7), Inches(0.38))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    num = slide.shapes.add_textbox(Inches(12.25), Inches(7.12), Inches(0.55), Inches(0.25))
    p = num.text_frame.paragraphs[0]
    p.text = str(number)
    p.font.name = "Arial"
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["muted"]
    p.alignment = PP_ALIGN.RIGHT


def add_card(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    set_fill(shape, fill or COLORS["card"])
    set_line(shape, line or COLORS["line"], 1)
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["ink"]
    return box


def add_bullets(slide, left, top, width, height, bullets, size=16):
    add_card(slide, left, top, width, height)
    box = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.18), Inches(width - 0.45), Inches(height - 0.3))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["ink"]
        p.space_after = Pt(7)


def add_badge(slide, left, top, width, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.38))
    set_fill(shape, color)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def add_image_card(slide, path, left, top, width, height, caption=None):
    add_card(slide, left, top, width, height)
    img_path = ROOT / path
    slide.shapes.add_picture(str(img_path), Inches(left + 0.12), Inches(top + 0.16), width=Inches(width - 0.24), height=Inches(height - (0.55 if caption else 0.28)))
    if caption:
        add_text(slide, left + 0.2, top + height - 0.35, width - 0.4, 0.22, caption, size=9, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def add_metric(slide, left, top, width, number, label, color):
    add_card(slide, left, top, width, 0.95, fill=RGBColor(239, 246, 255), line=RGBColor(191, 219, 254))
    add_text(slide, left + 0.15, top + 0.12, width - 0.3, 0.28, number, size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, left + 0.15, top + 0.48, width - 0.3, 0.26, label, size=10, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def make_slide(prs, title, number, bullets=None, image=None, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title, number)
    if image:
        add_bullets(slide, 0.65, 1.15, 5.65, 5.75, bullets or [], size=15)
        add_image_card(slide, image, 6.65, 1.15, 6.0, 5.75, caption)
    else:
        add_bullets(slide, 0.85, 1.15, 11.65, 5.75, bullets or [], size=17)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "Выпускная квалификационная работа", 1)
    add_text(slide, 0.85, 1.22, 8.2, 0.82, "Разработка и тестирование информационной системы для управления фитнес-клубом", size=28, bold=True, color=COLORS["dark"])
    add_text(slide, 0.9, 2.18, 7.0, 0.4, "Специальность 09.02.07 Информационные системы и программирование", size=16, color=COLORS["muted"])
    add_metric(slide, 0.9, 3.05, 3.55, "Flask", "веб-приложение", COLORS["blue"])
    add_metric(slide, 4.75, 3.05, 3.55, "SQLite", "единая база данных", COLORS["green"])
    add_metric(slide, 8.6, 3.05, 3.55, "pytest", "проверка сценариев", COLORS["amber"])
    add_card(slide, 0.9, 4.45, 11.25, 1.05, fill=RGBColor(240, 253, 244), line=RGBColor(187, 247, 208))
    add_text(slide, 1.15, 4.72, 10.7, 0.45, "Клиентская запись на тренировки, абонементы, административная панель, уведомления и аналитика", size=18, bold=True, color=RGBColor(20, 83, 45), align=PP_ALIGN.CENTER)

    slides = [
        ("Актуальность", ["Фитнес-клубу нужна единая цифровая среда для клиентов и администрации.", "Ручной учет приводит к ошибкам в расписании, абонементах и посещаемости.", "Клиент ожидает самостоятельную запись на занятия и прозрачный статус абонемента.", "Администрации нужны отчеты и быстрый доступ к актуальным данным."]),
        ("Цель и задачи", ["Цель: разработать и протестировать веб-систему управления фитнес-клубом.", "Проанализировать предметную область и сформировать требования.", "Спроектировать архитектуру приложения и базу данных.", "Реализовать пользовательские, административные и аналитические модули.", "Проверить основные сценарии автоматизированными тестами."]),
        ("Предметная область", ["Участники: клиент, тренер, администратор.", "Ключевые процессы: регистрация, расписание, запись, абонементы, посещаемость.", "Данные связаны между собой: клиент записывается на конкретную тренировку с конкретным тренером.", "Система должна поддерживать полный цикл обслуживания фитнес-клуба."]),
        ("Требования к системе", ["Регистрация, авторизация и разграничение ролей.", "Просмотр тренировок, тренеров и личного расписания.", "Запись на занятие и отмена будущей записи.", "Покупка, продление и заморозка абонемента.", "Управление пользователями, тренерами, расписанием и тарифами.", "Формирование аналитики и экспорт отчетов."]),
        ("Технологический стек", ["Python и Flask — серверная часть.", "SQLite и SQLAlchemy — хранение и обработка данных.", "Flask-Login — пользовательские сессии.", "Flask-WTF и WTForms — формы и валидация.", "Flask-Mail — уведомления.", "ReportLab, XlsxWriter и pytest — отчеты и тестирование."]),
        ("Архитектура проекта", ["Проект разделен на логические слои.", "controllers отвечают за маршруты и сценарии.", "models описывают сущности базы данных.", "forms проверяют ввод пользователя.", "services содержат уведомления, отчеты и вспомогательную логику."], "diagrams/project_structure.png", "Структура проекта Fitness Gym App"),
        ("База данных", ["User — пользователь и роль.", "Trainer — тренерский состав.", "WorkoutType — направление тренировки.", "Workout — занятие в расписании.", "Booking — запись пользователя.", "Abonement и UserAbonement — тариф и покупка."], "diagrams/er_diagram.png", "ER-диаграмма основной учетной модели"),
        ("Пользовательский сценарий", ["Регистрация и вход в систему.", "Просмотр расписания и фильтрация занятий.", "Просмотр карточки тренера.", "Запись на тренировку.", "Просмотр личного расписания.", "Работа с купленными абонементами."]),
        ("Административный сценарий", ["Вход администратора.", "Управление тренерами и типами тренировок.", "Формирование расписания.", "Управление пользователями и правами.", "Создание и редактирование абонементов.", "Переход к аналитике и отчетам."], "diagrams/Controllers Class Diagram.png", "Разделение маршрутов по blueprint-контроллерам"),
        ("Абонементы и уведомления", ["Каталог тарифов и детальная страница абонемента.", "Покупка, продление и заморозка.", "Расчет срока действия и остатка посещений.", "Уведомления о тренировках и состоянии абонемента.", "Отдельный сервисный слой для отчетов и сообщений."], "diagrams/services_class.png", "Сервисы отчетности и уведомлений"),
        ("Аналитика и отчеты", ["Популярность тренировок.", "Загрузка тренеров.", "Посещаемость клиентов.", "Выбор периода отчета.", "Экспорт в PDF и Excel.", "Поддержка управленческих решений на основе данных."]),
        ("Тестирование", ["pytest и тестовая база SQLite в памяти.", "Проверка главной страницы, входа и регистрации.", "Проверка пользовательских маршрутов.", "Проверка административного доступа.", "Проверка моделей User, Workout и Booking."], "diagrams/tests_class.png", "Структура автоматизированных тестов"),
        ("Результат разработки", ["Создана веб-ИС для управления фитнес-клубом.", "Снижена ручная нагрузка администратора.", "Централизованы расписание, клиенты и абонементы.", "Добавлены уведомления, аналитика и отчеты.", "Архитектура допускает дальнейшее развитие."]),
        ("Заключение", ["Цель выпускной квалификационной работы достигнута.", "Разработаны пользовательская и административная части.", "Реализованы абонементы, уведомления и отчетность.", "Проведено тестирование ключевых сценариев.", "Перспективы: онлайн-оплата, мобильная версия, кабинет тренера и расширенная аналитика."]),
    ]
    for idx, item in enumerate(slides, 2):
        title, bullets, *rest = item
        image = rest[0] if rest else None
        caption = rest[1] if len(rest) > 1 else None
        make_slide(prs, title, idx, bullets, image, caption)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
