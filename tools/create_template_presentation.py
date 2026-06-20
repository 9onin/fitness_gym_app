from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
TEMPLATE = DOCS / "+Шаблон презентации ВКР.pptx"
OUT = DOCS / "Черников_Никита_презентация_ВКР_по_шаблону.pptx"
ZIP_OUT = DOCS / "Черников_Никита_презентация_ВКР_по_шаблону.zip"
SCREENSHOTS = DOCS / "screenshots"
QR_PATH = DOCS / "github_qr_fitness_gym.png"

GITHUB_URL = "https://github.com/9onin/fitness_gym_app"
FULL_NAME = "Черников Никита Александрович"
SUPERVISOR = "Рыбаков Александр Сергеевич"
TOPIC = (
    "Отладка и тестирование информационной системы для учета клиентов "
    "фитнес-клуба с интеграцией мобильных приложений"
)

DARK = RGBColor(40, 40, 40)
ACCENT = RGBColor(31, 78, 121)


def clear_and_write(shape, paragraphs, font_size=Pt(24), bold_first=False, align=PP_ALIGN.LEFT):
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.08)
    text_frame.margin_right = Inches(0.08)
    text_frame.margin_top = Inches(0.04)
    text_frame.margin_bottom = Inches(0.04)

    if isinstance(paragraphs, str):
        paragraphs = [paragraphs]

    for idx, text in enumerate(paragraphs):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = text
        p.alignment = align
        p.space_after = Pt(4)
        p.line_spacing = 0.95
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = font_size
        run.font.color.rgb = DARK
        if idx == 0 and bold_first:
            run.font.bold = True


def write_bullets(shape, title, bullets, font_size=Pt(22)):
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.12)
    text_frame.margin_right = Inches(0.12)
    text_frame.margin_top = Inches(0.06)
    text_frame.margin_bottom = Inches(0.06)

    p = text_frame.paragraphs[0]
    p.text = title
    p.space_after = Pt(8)
    p.line_spacing = 0.95
    r = p.runs[0]
    r.font.name = "Arial"
    r.font.size = font_size
    r.font.bold = True
    r.font.color.rgb = ACCENT

    for bullet in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(4)
        p.line_spacing = 0.95
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(19)
        r.font.color.rgb = DARK


def replace_by_prefix(slide, prefix, paragraphs, font_size=Pt(22), title=None):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text.strip().startswith(prefix):
            if title:
                write_bullets(shape, title, paragraphs, font_size)
            else:
                clear_and_write(shape, paragraphs, font_size)
            return shape
    raise RuntimeError(f"Text block not found: {prefix}")


def remove_shape(shape):
    parent = shape._element.getparent()
    parent.remove(shape._element)


def add_picture_fit(slide, image_path, left, top, width, height):
    pic = slide.shapes.add_picture(str(image_path), left, top, width=width)
    if pic.height > height:
        ratio = height / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
    pic.left = int(left + (width - pic.width) / 2)
    pic.top = int(top + (height - pic.height) / 2)
    return pic


def create_qr():
    qr = qrcode.QRCode(box_size=12, border=2)
    qr.add_data(GITHUB_URL)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    img.save(QR_PATH)


def build():
    create_qr()
    prs = Presentation(TEMPLATE)

    # Slide 1
    replace_by_prefix(
        prs.slides[0],
        "Специальность:",
        [
            "Специальность: 09.02.07 Информационные системы и программирование",
            "Квалификация: Специалист по информационным системам",
            "Курс IV, группа ИПО-41.22",
            f"Обучающийся: {FULL_NAME}",
            f"Руководитель: {SUPERVISOR}",
        ],
        Pt(19),
    )
    replace_by_prefix(prs.slides[0], "«Разработка", f"«{TOPIC}»", Pt(27))

    replacements = {
        2: (
            "Цель работы",
            [
                "разработать и проверить информационную систему учета клиентов фитнес-клуба;",
                "изучить процессы записи, абонементов, расписания и администрирования;",
                "спроектировать структуру приложения, базу данных и пользовательские сценарии;",
                "реализовать клиентскую и административную части веб-приложения;",
                "провести отладку и тестирование ключевых функций системы.",
            ],
        ),
        3: (
            "Проблема автоматизации",
            [
                "ручной учет клиентов, записей и абонементов приводит к ошибкам и потере времени;",
                "расписание, тренеры, тренировки и абонементы должны храниться согласованно;",
                "клиенту нужен быстрый доступ к занятиям, а администратору - единая панель управления;",
                "собственная система позволяет адаптировать функции под реальные процессы клуба.",
            ],
        ),
        4: (
            "Предметная область",
            [
                "участники: посетитель, администратор, тренер и информационная система;",
                "основные данные: пользователи, тренеры, типы тренировок, расписание, записи, абонементы;",
                "ключевые процессы: регистрация, авторизация, выбор занятия, бронирование, управление тарифами;",
                "по итогам анализа сформированы функциональные требования к системе.",
            ],
        ),
        5: (
            "Выбор подхода",
            [
                "готовые сервисы часто избыточны, платны или плохо адаптируются под учебный проект;",
                "для небольшого фитнес-клуба важны простое развертывание и понятная структура данных;",
                "выбрана самостоятельная разработка веб-приложения с модульной архитектурой;",
                "это дает контроль над логикой учета, интерфейсом и дальнейшим расширением.",
            ],
        ),
        6: (
            "Архитектура решения",
            [
                "приложение построено на Flask и разделено на логические blueprint-модули;",
                "выделены публичная часть, авторизация, личный кабинет, абонементы, админ-панель и аналитика;",
                "модели SQLAlchemy описывают основные сущности и связи базы данных;",
                "разделение маршрутов и шаблонов упрощает сопровождение и тестирование.",
            ],
        ),
        7: (
            "Технологический стек",
            [
                "Python и Flask - серверная логика и маршрутизация веб-приложения;",
                "SQLAlchemy и SQLite - хранение данных и работа с моделями;",
                "Flask-Login и WTForms - авторизация, сессии и проверка форм;",
                "HTML, CSS, Bootstrap/Tailwind-подходы - интерфейс пользователя и администратора;",
                "pytest - автоматизированная проверка основных сценариев.",
            ],
        ),
        8: (
            "База данных и модули",
            [
                "основные сущности: User, Trainer, WorkoutType, Workout, Booking, Membership;",
                "связи БД поддерживают запись клиента на тренировку и учет абонементов;",
                "модули системы покрывают регистрацию, расписание, тренеров, записи, абонементы и отчеты;",
                "административный модуль управляет справочниками и клиентской базой.",
            ],
        ),
        9: (
            "Функционал и проверка",
            [
                "посетитель просматривает главную страницу, тренировки, абонементы и оформляет запись;",
                "администратор управляет пользователями, тренерами, расписанием, тарифами и отчетами;",
                "тестирование охватывает публичные страницы, авторизацию, пользовательские маршруты, админ-доступ и модели;",
                "отладка выполнялась по маршрутам, формам, шаблонам, связям моделей и правам доступа.",
            ],
        ),
        10: (
            "Итог разработки",
            [
                "создано работающее веб-приложение для учета клиентов фитнес-клуба;",
                "реализованы регистрация, авторизация, расписание, записи, абонементы, уведомления и аналитика;",
                "подготовлены диаграммы архитектуры, БД, классов, прецедентов и тестовой структуры;",
                "результат соответствует цели ВКР и подтвержден тестированием основных сценариев.",
            ],
        ),
        11: (
            "Практическая значимость",
            [
                "система снижает объем ручной работы администратора;",
                "клиент получает самостоятельный доступ к расписанию и абонементам;",
                "данные о пользователях, тренировках и посещениях хранятся централизованно;",
                "личный вклад включает анализ, проектирование, реализацию интерфейсов, БД и тестов.",
            ],
        ),
        12: (
            "Выводы и развитие",
            [
                "цель работы достигнута, поставленные задачи выполнены;",
                "разработанная система может служить основой для автоматизации фитнес-клуба;",
                "перспективы: онлайн-оплата, мобильное приложение, кабинет тренера, push-уведомления;",
                "также возможно расширение аналитики, безопасности и интеграций с внешними сервисами.",
            ],
        ),
    }

    red_prefixes = {
        2: "На этом слайде",
        3: "Здесь раскрывается",
        4: "В аналитической",
        5: "На защите",
        6: "Практическую",
        7: "На этом слайде",
        8: "Этот слайд",
        9: "Здесь нужно",
        10: "Этот слайд",
        11: "На этом слайде",
        12: "В заключении",
    }

    for slide_no, (title, bullets) in replacements.items():
        slide = prs.slides[slide_no - 1]
        replace_by_prefix(slide, red_prefixes[slide_no], bullets, Pt(22), title)

    # Slide 9 screenshots
    slide = prs.slides[8]
    text_shape = list(slide.shapes)[2]
    text_shape.top = Inches(0.95)
    text_shape.height = Inches(1.78)
    screenshots = [
        SCREENSHOTS / "homepage.png",
        SCREENSHOTS / "visitor_workouts.png",
        SCREENSHOTS / "admin_dashboard.png",
    ]
    positions = [
        (Inches(0.45), Inches(3.05), Inches(3.85), Inches(2.45)),
        (Inches(4.75), Inches(3.05), Inches(3.85), Inches(2.45)),
        (Inches(9.05), Inches(3.05), Inches(3.85), Inches(2.45)),
    ]
    for img, pos in zip(screenshots, positions):
        if img.exists():
            add_picture_fit(slide, img, *pos)

    # Slide 13
    replace_by_prefix(
        prs.slides[12],
        "На финальном",
        "Работа завершена. Основные результаты представлены, проект готов к демонстрации и обсуждению.",
        Pt(20),
    )
    replace_by_prefix(
        prs.slides[12],
        "Репозиторий GitHub:",
        f"Репозиторий GitHub: {GITHUB_URL}",
        Pt(19),
    )
    replace_by_prefix(prs.slides[12], "Вставить QR", "QR-код репозитория", Pt(19))
    prs.slides[12].shapes.add_picture(str(QR_PATH), Inches(5.55), Inches(4.55), width=Inches(1.35))

    prs.save(OUT)
    if ZIP_OUT.exists():
        ZIP_OUT.unlink()
    with ZipFile(ZIP_OUT, "w", ZIP_DEFLATED) as z:
        z.write(OUT, OUT.name)
    return OUT, ZIP_OUT


if __name__ == "__main__":
    out, zip_out = build()
    print(out)
    print(zip_out)
